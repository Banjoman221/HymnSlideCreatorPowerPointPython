import imp
from itertools import count
from pydoc import pathdirs, visiblename
import csv
import os, sys
from os import startfile
import win32com.client as win32
import msvcrt

# import Presentation class
# from pptx library
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

application_path = os.path.dirname(sys.executable)

# Getting main path of this folder
mainPath = os.path.dirname(__file__)
# Getting CSV file
hymn = os.path.join(mainPath, "hymnlist.csv")

# Accessing CSV file and adding to an array to be accessed later
with open(hymn, newline="") as csvfile:
    rows = csv.reader(csvfile)
    data = []
    for row in rows:
        data.append(row)

while True:
    print("At any time q will exit the program")
    countar = []
    noh = []
    number = 0
    hhn = 0
    p = ""
    print("How many Hymns 1-5:")
    howManyHymns = msvcrt.getwch()
    qt = "Q"

    if howManyHymns.capitalize() == qt:
        quit()

    # Determing how many hymns
    while True:
        try:
            if howManyHymns.isnumeric() == True and int(howManyHymns) <= 5:
                hhn = howManyHymns
                break
            elif howManyHymns.capitalize() == qt:
                quit()
            else:
                print("How many Hymns 1-5:")
                howManyHymns = msvcrt.getwch()
        except:
            print("How many Hymns 1-5:")
            howManyHymns = msvcrt.getwch()
            if howManyHymns.capitalize() == qt:
                quit()
    print(hhn)

    # Getting Hymn numbers
    while len(countar) != int(hhn):
        # Getting input for which number of hymn
        naoh = input("Enter a hymn number between 1 and 479: ")
        # Making sure the input is a number and less than 479. If q, quit script
        if naoh.isnumeric() and int(naoh) <= 479 and int(naoh) != 0:
            # Getting the Name of the Hymn from CSV data array
            name = data[int(naoh) - 1][0]
            # Adding Hymn name to an array - countar
            if naoh not in noh:
                countar.append(name)
                # Adding Hymn number to an array - noh
                noh.append(naoh)
                print(name)
            else:
                print("Slide already in use please try again")
        elif naoh.isnumeric() and int(naoh) >= 479:
            naoh = input("Enter a hymn number between 1 and 479: ")
        elif naoh.capitalize() == qt:
            quit()

    # Creating presentation object
    root = Presentation()
    # Setting the width and heighth of a slide
    root.slide_width = Inches(16)
    root.slide_height = Inches(9)

    print(countar)

    # Creating a slide from the Hymn numbers given
    if len(countar) == int(hhn):
        # Getting the name of the slide
        for name in countar:
            number += 1
            print(number)
            # Creating slide layout
            first_slide_layout = root.slide_layouts[6]
            # Adding a slide
            slide = root.slides.add_slide(first_slide_layout)
            # Getting slide background
            path = os.path.join(mainPath, "bg.jpg")

            # Making sure left and top inches = 0 for background image use
            left = top = Inches(0)
            # Creating slide background shape and appling background
            pic = slide.shapes.add_picture(
                path, left, top, width=root.slide_width, height=root.slide_height
            )
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            class Shape:
                def __init__(
                    self,
                    left,
                    top,
                    width,
                    height,
                    tname,
                    fname,
                    fsize,
                    fbold,
                    falignment,
                    fcr,
                ):
                    # Creating a shape for Hymn Title
                    self.left = Inches(left)
                    self.top = Inches(top)
                    self.width = Inches(width)
                    self.height = Inches(height)
                    txBox = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        self.left,
                        self.top,
                        self.width,
                        self.height,
                    )
                    # Making shape match background
                    txBox.fill.background()
                    txBox.line.fill.background()
                    # Creating a text frame
                    tf = txBox.text_frame
                    # Adding Text
                    self.tname = tname
                    tp = tf.paragraphs[0]
                    run = tp.add_run()
                    run.text = str(self.tname)
                    # Customizing text
                    self.fname = fname
                    self.fsize = fsize
                    self.fbold = fbold
                    self.falignment = falignment
                    self.fcr = fcr

                    tp.font.name = fname
                    tp.font.size = fsize
                    tp.font.bold = fbold
                    tp.alignment = falignment
                    tp.font.color.rgb = fcr

            Shape(
                0,
                0.5,
                16,
                3.47,
                name,
                "ALGERIAN",
                Pt(115),
                True,
                PP_ALIGN.CENTER,
                RGBColor(0, 0, 0),
            )

            Shape(
                0,
                5,
                16,
                3.47,
                noh[number - 1],
                "ALGERIAN",
                Pt(190),
                True,
                PP_ALIGN.CENTER,
                RGBColor(0, 0, 0),
            )

    # Grabing the powerpoint application
    powerpoint = win32.gencache.EnsureDispatch("PowerPoint.Application")
    # Quiting the powerpoint appliction if running
    powerpoint.Quit()

    pptout = os.path.join(mainPath, "Output.pptx")
    # Saving file
    root.save(pptout)
    # Starting file
    startfile(pptout)

    print("done")
